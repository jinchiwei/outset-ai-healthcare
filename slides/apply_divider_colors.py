"""Post-build pass: repaint section-divider slides full-bleed with their cycling
brand accent (turquoise -> deeppink -> amber -> blueviolet), with contrast-correct
text and structural marks.

The renderer draws dividers on navy by default (themes never paint with the brand-4).
This pass reads each divider's accent_hex from the layout sidecar, maps it to its
rendered slide, repaints the full-bleed background to that accent, and flips the
title / eyebrow / bars / footer to ink-or-white for contrast.

Run after build + before/after apply_logo:
    python slides/apply_divider_colors.py slides/build/day1.pptx slides/day1.md.layout.json
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
import json

LIGHT = {"#40E0D0", "#F0C840"}  # turquoise, amber -> ink text
INK = RGBColor(0x14, 0x14, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EMU_PER_IN = 914400

# Mirror render.py: section accent cycles the brand-4 by first-seen section order
# (NOT the per-entry accent_hex in the sidecar, which can be stale).
SECTION_CYCLE = ["#40E0D0", "#FF1493", "#F0C840", "#8A2BE2"]  # turq, pink, amber, violet


def _rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _contrast(accent_hex):
    return INK if accent_hex.upper() in {c.upper() for c in LIGHT} else WHITE


def _set_text_color(shape, color):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = color


def apply(pptx_path: Path, sidecar_path: Path):
    plan = json.loads(sidecar_path.read_text())
    prs = Presentation(str(pptx_path))
    n_slides = len(prs.slides)

    # Replicate render.py's first-seen section cycling to get each divider's accent.
    seen = []

    def section_accent(label):
        if label and label not in seen:
            seen.append(label)
        i = seen.index(label) if label in seen else 0
        return SECTION_CYCLE[i % len(SECTION_CYCLE)]

    # plan slides render to prs.slides[1 + position] (slide[0] = cover)
    recolored = 0
    for pos, entry in enumerate(plan["slides"]):
        label = entry.get("params", {}).get("label") or entry.get("params", {}).get("section_label", "")
        if entry["kind"] != "section-divider":
            # still advance first-seen tracking on content slides' section_label
            section_accent(label)
            continue
        accent_hex = section_accent(label)
        accent = _rgb(accent_hex)
        text_color = _contrast(accent_hex)
        idx = 1 + pos
        if idx >= n_slides:
            continue
        slide = prs.slides[idx]
        for sh in slide.shapes:
            is_full_bleed = (sh.width and sh.height
                             and sh.width > 13 * EMU_PER_IN and sh.height > 7 * EMU_PER_IN)
            if sh.has_text_frame and sh.text_frame.text.strip():
                _set_text_color(sh, text_color)
            try:
                if is_full_bleed:
                    sh.fill.solid(); sh.fill.fore_color.rgb = accent
                    sh.line.fill.background()
                elif sh.shape_type is not None and not sh.has_text_frame:
                    # thin accent marks (colorblock, top bar, hairline) -> contrast
                    sh.fill.solid(); sh.fill.fore_color.rgb = text_color
                    sh.line.fill.background()
            except (AttributeError, TypeError, ValueError):
                pass
        recolored += 1
    prs.save(str(pptx_path))
    print(f"recolored {recolored} section dividers full-bleed")


if __name__ == "__main__":
    pptx = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("slides/build/day1.pptx")
    side = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("slides/day1.md.layout.json")
    apply(pptx, side)
