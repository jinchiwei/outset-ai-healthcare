"""Post-build pass: co-author treatment for the Outset decks.

1. Recolor the cover byline (the names) to matplotlib orange.
2. Append a custom closing/thank-you slide: dark, with just the two English
   names in matplotlib orange (Curtis first), no email / CJK / identity logos.

Run after the deck is built (and after apply_logo):
    python slides/apply_authors.py slides/build/day1.pptx
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

MPL_ORANGE = RGBColor(0xFF, 0x7F, 0x0E)   # matplotlib default 'tab:orange' / C1
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TURQUOISE = RGBColor(0x40, 0xE0, 0xD0)
DEEPPINK = RGBColor(0xFF, 0x14, 0x93)
DARK = RGBColor(0x14, 0x14, 0x1C)
MUTED = RGBColor(0x9A, 0x9A, 0xA8)
MONO = "Geist Mono"

# co-authors, Curtis first
NAMES = ["Curtis Chambers", "Jinchi Wei"]
SLIDE_W, SLIDE_H = 13.333, 7.5


def _blank_layout(prs):
    for lay in prs.slide_layouts:
        if len(lay.placeholders) == 0:
            return lay
    return prs.slide_layouts[-1]


def _text(slide, text, *, left, top, width, height, size, color, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = MONO
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _rect(slide, *, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def recolor_cover_byline(prs):
    cover = prs.slides[0]
    for sh in cover.shapes:
        if sh.has_text_frame and "Jinchi" in sh.text_frame.text and "Curtis" in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = MPL_ORANGE
            return True
    return False


def add_closing(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _rect(slide, left=0, top=0, width=SLIDE_W, height=SLIDE_H, color=DARK)
    _rect(slide, left=0, top=0, width=0.8, height=SLIDE_H, color=TURQUOISE)
    _rect(slide, left=0.8, top=0, width=0.25, height=SLIDE_H, color=DEEPPINK)

    _text(slide, "THANK YOU", left=1.3, top=1.7, width=11.0, height=1.4,
          size=58, color=WHITE)
    _text(slide, NAMES[0], left=1.3, top=3.5, width=11.0, height=0.7,
          size=30, color=MPL_ORANGE)
    _text(slide, NAMES[1], left=1.3, top=4.25, width=11.0, height=0.7,
          size=30, color=MPL_ORANGE)
    _text(slide, "AI in Healthcare  ·  Outset", left=1.3, top=5.5, width=11.0, height=0.5,
          size=15, color=MUTED, bold=False)
    return slide


def apply(pptx_path: Path):
    prs = Presentation(str(pptx_path))
    ok = recolor_cover_byline(prs)
    add_closing(prs)
    prs.save(str(pptx_path))
    print(f"co-authors applied to {pptx_path.name} (cover byline orange: {ok})")


if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("slides/build/day1.pptx")
    apply(target)
