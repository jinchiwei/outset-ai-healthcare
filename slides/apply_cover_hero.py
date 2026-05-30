"""Post-build pass: place a framed real medical image on the right of the cover.

The cover text sits in the left ~55%; this drops a clean square hero image in the
open right area, with a thin accent frame. Subtle, not decorative clutter.

    python slides/apply_cover_hero.py slides/build/day1.pptx fundus_dr.jpg 40E0D0
"""
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
REALIMG = ROOT / "slides/realimg"
SLIDE_W, SLIDE_H = 13.333, 7.5


def apply(pptx_path: Path, img_name: str, accent_hex: str):
    src = REALIMG / img_name
    im = Image.open(src)
    s = min(im.size)
    im = im.crop(((im.width - s) // 2, (im.height - s) // 2,
                  (im.width + s) // 2, (im.height + s) // 2))
    sq = REALIMG / f".cover_{img_name}.png"
    im.save(sq)

    # Right side, below the title's first line and above the logo. Keep cover
    # subtitles short enough to end left of this image (~x10).
    size = 2.4                      # inches, square
    left = SLIDE_W - 0.62 - size
    top = 3.5
    accent = RGBColor.from_string(accent_hex)

    prs = Presentation(str(pptx_path))
    cover = prs.slides[0]
    # accent frame slightly larger than the image
    pad = 0.06
    fr = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left - pad), Inches(top - pad),
                                Inches(size + 2 * pad), Inches(size + 2 * pad))
    fr.fill.solid(); fr.fill.fore_color.rgb = accent
    fr.line.fill.background(); fr.shadow.inherit = False
    cover.shapes.add_picture(str(sq), Inches(left), Inches(top),
                             width=Inches(size), height=Inches(size))
    prs.save(str(pptx_path))
    print(f"cover hero ({img_name}) applied to {pptx_path.name}")


if __name__ == "__main__":
    pptx = Path(sys.argv[1])
    img = sys.argv[2] if len(sys.argv) > 2 else "fundus_dr.jpg"
    accent = sys.argv[3] if len(sys.argv) > 3 else "40E0D0"
    apply(pptx, img, accent)
