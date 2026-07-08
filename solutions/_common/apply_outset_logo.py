"""Post-build: drop the Outset logo (white, for the dark cover) on a deck's title slide.

Usage:  python solutions/_common/apply_outset_logo.py <deck.pptx>
Uses the repo asset slides/assets/outset-logo-white.png (generated from outset-logo.png if absent).
"""
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "slides/assets/outset-logo.png"
WHITE = ROOT / "slides/assets/outset-logo-white.png"
SLIDE_W, SLIDE_H = 13.333, 7.5
LOGO_H, MARGIN = 0.5, 0.55


def _ensure_white():
    if WHITE.exists():
        return
    im = Image.open(SRC).convert("RGBA")
    im.putdata([(255, 255, 255, a) for (_r, _g, _b, a) in im.getdata()])
    im.save(WHITE)


def apply(pptx_path):
    _ensure_white()
    aspect = Image.open(WHITE).width / Image.open(WHITE).height
    w = LOGO_H * aspect
    prs = Presentation(str(pptx_path))
    prs.slides[0].shapes.add_picture(str(WHITE), Inches(SLIDE_W - MARGIN - w),
                                     Inches(SLIDE_H - MARGIN - LOGO_H), height=Inches(LOGO_H))
    prs.save(str(pptx_path))
    print(f"applied Outset logo to cover of {Path(pptx_path).name}")


if __name__ == "__main__":
    apply(sys.argv[1])
