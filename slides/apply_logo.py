"""Post-build pass: drop the Outset wordmark on the cover slide (lower-right).

The cover is dark navy, so we use a white-recolored version of the wordmark.
Run after each build (the renderer overwrites the pptx, so the logo must be
re-applied):  python slides/apply_logo.py slides/build/day1.pptx
"""
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
SRC_LOGO = ROOT / "slides/assets/outset-logo.png"
WHITE_LOGO = ROOT / "slides/assets/outset-logo-white.png"

SLIDE_W = 13.333
SLIDE_H = 7.5
LOGO_H = 0.42          # inches
MARGIN = 0.55          # from right + bottom edges


def ensure_white_logo():
    if WHITE_LOGO.exists():
        return
    im = Image.open(SRC_LOGO).convert("RGBA")
    px = im.getdata()
    out = [(255, 255, 255, a) for (_r, _g, _b, a) in px]
    im.putdata(out)
    im.save(WHITE_LOGO)
    print("wrote white logo")


def apply(pptx_path: Path):
    ensure_white_logo()
    im = Image.open(WHITE_LOGO)
    aspect = im.width / im.height
    w = LOGO_H * aspect
    left = SLIDE_W - MARGIN - w
    top = SLIDE_H - MARGIN - LOGO_H
    prs = Presentation(str(pptx_path))
    prs.slides[0].shapes.add_picture(str(WHITE_LOGO), Inches(left), Inches(top),
                                     height=Inches(LOGO_H))
    prs.save(str(pptx_path))
    print(f"applied Outset logo to cover of {pptx_path.name} ({w:.2f}in wide)")


if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "slides/build/day1.pptx"
    apply(target)
