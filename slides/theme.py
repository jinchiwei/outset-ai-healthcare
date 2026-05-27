"""Brand theme for the Outset AI in Healthcare slide decks.

Colors and fonts mirror Jinchi's personal branding: Geist / Geist Mono, with the
canonical palette turquoise / deeppink / amber / blueviolet. All decks are 16:9.

Public helpers:
    new_deck()       -> a 16:9 Presentation
    title_slide(...) -> opening slide with accent bar
    content_slide(...) -> title + underline accent + bullets
    section_divider(...) -> full-bleed dark section break
    save(prs, path)  -> write the .pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

# Brand palette (see memory/user_branding.md). Canonical priority colors.
TURQUOISE = RGBColor(0x40, 0xE0, 0xD0)
DEEPPINK = RGBColor(0xFF, 0x14, 0x93)
AMBER = RGBColor(0xF0, 0xC8, 0x40)
BLUEVIOLET = RGBColor(0x8A, 0x2B, 0xE2)
INK = RGBColor(0x10, 0x10, 0x10)
PAPER = RGBColor(0xFA, 0xFA, 0xFA)
MUTED = RGBColor(0x66, 0x66, 0x66)

SANS = "Geist"
MONO = "Geist Mono"

# Canonical PowerPoint 16:9 widescreen dimensions, in EMU (914400 EMU/inch).
SLIDE_W = Emu(12192000)  # 13.333 in
SLIDE_H = Emu(6858000)   # 7.5 in

_BLANK = 6  # blank layout index in the default template


def new_deck() -> Presentation:
    """Return a fresh 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[_BLANK])


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def title_slide(prs, *, title, subtitle="", accent=TURQUOISE):
    """Opening slide: vertical accent bar, large title, muted subtitle."""
    slide = _add_blank(prs)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), Inches(0.4), Inches(6.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tx = _textbox(slide, Inches(1.0), Inches(2.5), Inches(11.5), Inches(2.0))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = SANS
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = INK

    if subtitle:
        sx = _textbox(slide, Inches(1.0), Inches(4.6), Inches(11.5), Inches(1.0))
        sp = sx.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = SANS
        sp.font.size = Pt(24)
        sp.font.color.rgb = MUTED

    return slide


def content_slide(prs, *, title, bullets, accent=TURQUOISE):
    """Title with an underline accent, then a list of bullets."""
    slide = _add_blank(prs)

    tx = _textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = SANS
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INK

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(1.6), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    bx = _textbox(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
    tf = bx.text_frame
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.font.name = SANS
        para.font.size = Pt(22)
        para.font.color.rgb = INK
        para.space_after = Pt(10)

    return slide


def section_divider(prs, *, label, accent=DEEPPINK):
    """Full-bleed dark slide with a large accent-colored label."""
    slide = _add_blank(prs)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()

    tx = _textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = label
    p.font.name = SANS
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = accent

    return slide


def code_slide(prs, *, title, code, accent=AMBER):
    """Title + a monospace code block (for 'image.shape == (224,224,3)' style slides)."""
    slide = _add_blank(prs)

    tx = _textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = SANS
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INK

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(1.6), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    cx = _textbox(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(5.0))
    tf = cx.text_frame
    for i, line in enumerate(code.splitlines()):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.name = MONO
        para.font.size = Pt(20)
        para.font.color.rgb = INK

    return slide


def save(prs, path) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
